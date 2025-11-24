from django.core.files.storage import default_storage
from django.contrib.staticfiles.finders import find
import importlib.resources
from pptx import Presentation
from pptx.util import Inches, Pt
from django.utils.timezone import now


def url_fetcher(url, timeout=5, ssl_context=None):
    from weasyprint import default_url_fetcher

    # handle local media
    if url.startswith("local:"):
        path = url[6:]
        file_obj = default_storage.open(path, "rb")
        return {
            "file_obj": file_obj,
        }
    if url.startswith("static:"):
        path = url[7:]
        abs_path = find(path)
        return {"file_obj": open(abs_path, "rb")}

    return default_url_fetcher(url, timeout=timeout, ssl_context=ssl_context)


def render_anniversary_report_pptx(target, operation_name, month, incidents, actions, gaps):
    prs = Presentation()

    title_slide_layout = prs.slide_layouts[0]
    title_and_content_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Defect Elimination\n\nSignificant Incident Anniversaries"
    subtitle.text = f"{operation_name} - {month}"

    slide = prs.slides.add_slide(title_and_content_slide_layout)
    slide.shapes.title.text = "Significant Reliability Incidents Identified"
    body = slide.placeholders[1]

    headers = [
        "Section",
        "Date of Incident",
        "Equipment",
        "Description",
        "Solutions Tracked",
        "Solutions Verified",
    ]

    frame = slide.shapes.add_table(
        rows=2,
        cols=len(headers),
        left=body.left,
        top=body.top,
        width=body.width,
        height=body.height,
    )

    table = frame.table

    for ix, header in enumerate(headers):
        table.cell(0, ix).text = header

    for cell in table.iter_cells():
        cell.text_frame.margin_left = Pt(4)
        cell.text_frame.margin_right = Pt(4)
        cell.text_frame.margin_top = Pt(2)
        cell.text_frame.margin_bottom = Pt(2)
        cell.vertical_anchor = 1  # MSO_VERTICAL_ANCHOR.TOP (or use MIDDLE=2, BOTTOM=3)

    for ix, incident in enumerate(incidents, start=1):

        solutions = list(incident.solutions.all())
        solutions_tracked = len(solutions) > 0
        solutions_verified = all(solution.date_verified for solution in solutions)

        table.cell(ix, 0).text = incident.section.name
        table.cell(ix, 1).text = incident.time_start.strftime("%Y-%m-%d")
        table.cell(ix, 2).text = incident.equipment.name
        table.cell(ix, 3).text = incident.short_description
        table.cell(ix, 4).text = "yes" if solutions_tracked else "no" # solutions tracked
        table.cell(ix, 5).text = "yes" if solutions_verified else "no" # solutions verified

    slide = prs.slides.add_slide(title_and_content_slide_layout)
    slide.shapes.title.text = "Close-out Actions Identified"
    body = slide.placeholders[1]

    headers = [
        "Section",
        "Incident",
        "Proposed Solutions",
        "Progress Status",
        "Solution Verified",
        "Verification Comment"
    ]

    frame = slide.shapes.add_table(
        rows=2,
        cols=len(headers),
        left=body.left,
        top=body.top,
        width=body.width,
        height=body.height,
    )
    table = frame.table

    for ix, header in enumerate(headers):
        table.cell(0, ix).text = header

    for cell in table.iter_cells():
        cell.text_frame.margin_left = Pt(4)
        cell.text_frame.margin_right = Pt(4)
        cell.text_frame.margin_top = Pt(2)
        cell.text_frame.margin_bottom = Pt(2)
        cell.vertical_anchor = 1  # MSO_VERTICAL_ANCHOR.TOP (or use MIDDLE=2, BOTTOM=3)

    solutions = []
    for incident in incidents:

        for solution in incident.solutions.all():
            solutions.append({
                "section": incident.section.name,
                "incident": incident.short_description,
                "proposed_solution": solution.description,
                "progress_status": solution.status,
                "solution_verified": "yes" if solution.date_verified else "no",
                "verification_comment": solution.verification_comment,
            })

        solutions += list(incident.solutions.all())

    for ix, solution in enumerate(solutions, start=1):
        table.cell(ix, 0).text = solution["section"]
        table.cell(ix, 1).text = solution["incident"]
        table.cell(ix, 2).text = solution["proposed_solution"]
        table.cell(ix, 3).text = solution["progress_status"]
        table.cell(ix, 4).text = solution["solution_verified"]
        table.cell(ix, 5).text = solution["verification_comment"]

    slide = prs.slides.add_slide(title_and_content_slide_layout)
    slide.shapes.title.text = "Gaps and Close-Outs Identified"
    body = slide.placeholders[1]

    headers = [
        "Section",
        "Incident",
        "Potential Gaps Identified",
        "Potential Close-out Identified",
        "Commitment Signature",
    ]

    frame = slide.shapes.add_table(
        rows=2,
        cols=len(headers),
        left=body.left,
        top=body.top,
        width=body.width,
        height=body.height,
    )
    table = frame.table

    for ix, header in enumerate(headers):
        table.cell(0, ix).text = header

    for cell in table.iter_cells():
        cell.text_frame.margin_left = Pt(4)
        cell.text_frame.margin_right = Pt(4)
        cell.text_frame.margin_top = Pt(2)
        cell.text_frame.margin_bottom = Pt(2)
        cell.vertical_anchor = 1  # MSO_VERTICAL_ANCHOR.TOP (or use MIDDLE=2, BOTTOM=3)

    for ix, incident in enumerate(incidents, start=1):
        table.cell(ix, 0).text = incident.section.name
        table.cell(ix, 1).text = incident.short_description

    prs.save(target)
